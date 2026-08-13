"""
SENTINEL — Master PowerPoint (.pptx) Pitch Deck Generator
Structured specifically for Sivabalan T (Sri Sairam Engineering College):
1. Introduction to Myself & My Project
2. Core Problem Statement
3. Existing Solutions & Critical Drawbacks
4. The SENTINEL Solution & Innovation Matrix
5. Full End-to-End System Architecture & Flowchart Mapping
6. Architectural Provenance & Open-Source Adaptations
7. Zero-Trust Sanitizer & Dual-View Interface
8. 3-Tier System-Level MoE AI Routing Engine
9. AST Code Sandbox Guard & Active Defense
10. Courtroom PDF Briefs & Audit Trail
11. Enterprise Impact & System Verification Victory
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

def create_master_structured_pptx(output_path=r"C:\Users\siva2\Desktop\SENTINEL_Enterprise_Pitch_Deck.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Dark Enterprise Palette
    BG_NAVY = RGBColor(9, 13, 22)
    TEXT_WHITE = RGBColor(248, 250, 252)
    TEXT_MUTED = RGBColor(148, 163, 184)
    ACCENT_CYAN = RGBColor(6, 182, 212)
    ACCENT_BLUE = RGBColor(96, 165, 250)
    ACCENT_YELLOW = RGBColor(234, 179, 8)

    slides_content = [
        # SLIDE 1: INTRODUCTION TO MYSELF & MY PROJECT
        {
            "slide_title": "SENTINEL",
            "subtitle": "Autonomous AI SOC Triage & Zero-Trust Privacy Platform",
            "takeaway": "ENGINEERING INNOVATION: Next-Generation Privacy-Preserving Autonomous Cyber Defense Platform",
            "cards": [
                {"icon": "🎓", "title": "Presenter & Affiliation", "text": "<b>Sivabalan T</b>\nDepartment of Computer Science & Engineering (2nd Year)\nSri Sairam Engineering College"},
                {"icon": "🛡️", "title": "My Project Mission", "text": "Building an autonomous, privacy-preserving AI SOC analyst from first principles to solve alert fatigue and prevent data leakage."},
                {"icon": "🚀", "title": "Core Breakthroughs", "text": "Zero-Trust Data Sanitizer + 3-Tier MoE AI Router + AST Code Sandbox Guard + Courtroom PDF Briefs"}
            ],
            "notes": "Good morning respected judges. I am Sivabalan T, 2nd year Computer Science Engineering student from Sri Sairam Engineering College. Today I present my project, SENTINEL—an autonomous AI SOC triage and zero-trust privacy platform."
        },
        # SLIDE 2: CORE PROBLEM STATEMENT
        {
            "takeaway": "THE CRISIS: Security Operations Centers face catastrophic alert fatigue while zero-day attacks slip through.",
            "slide_title": "🚨 The Core Problem Statement",
            "subtitle": "Why Modern SOC Operations Are Overwhelmed and Vulnerable",
            "cards": [
                {"icon": "⚡", "title": "5,000+ Daily SIEM Alerts", "text": "SOC teams handle over 5,000 raw logs daily. Manual log parsing takes 30-45 minutes per alert, leading to 70%+ ignored alerts."},
                {"icon": "⏳", "title": "Critical Time-to-Detect Delay", "text": "Human triage latency allows ransomware and advanced persistent threats (APTs) to dwell undetected for weeks."},
                {"icon": "🎯", "title": "Lack of Contextual Correlation", "text": "Isolated SIEM logs fail to connect lateral movement across users, hosts, processes, and network domains."}
            ],
            "notes": "Let us look at the core problem: SOC teams face over 5,000 alerts daily. Manual triage takes 30 to 45 minutes per alert, causing critical ransomware attacks to slip through."
        },
        # SLIDE 3: EXISTING SOLUTIONS & THEIR DRAWBACKS
        {
            "takeaway": "CURRENT LIMITATIONS: Commercial AI wrappers leak sensitive police data, incur high costs, and risk system takeover.",
            "slide_title": "⚠️ Existing Solutions & Critical Drawbacks",
            "subtitle": "Why Current Commercial AI Wrappers and Legacy Tools Fail",
            "cards": [
                {"icon": "🔓", "title": "Drawback 1: Cloud PII Leakage", "text": "Tools like AiSOC send raw police emails, passwords, and internal IPs to public cloud LLMs, violating DPDP Act 2023 & GDPR."},
                {"icon": "💸", "title": "Drawback 2: High API Token Costs", "text": "Paying commercial cloud APIs per token on raw multi-megabyte log streams costs thousands of dollars monthly."},
                {"icon": "💥", "title": "Drawback 3: Unchecked Code Execution", "text": "Executing AI-generated scripts without AST verification allows prompt injection command execution (os.system)."}
            ],
            "notes": "Existing commercial AI wrappers have three major flaws: First, they leak raw police PII to public clouds; second, they burn thousands of dollars in cloud token fees; and third, they execute AI code blindly without security checks."
        },
        # SLIDE 4: THE SENTINEL SOLUTION & INNOVATION MATRIX
        {
            "takeaway": "MY SOLUTION: A 4-pillar privacy-first AI architecture that operates locally for $0 cost with AST code safety.",
            "slide_title": "💡 The SENTINEL Solution & Innovation Matrix",
            "subtitle": "Engineered Solutions Addressing Every Existing Drawback",
            "cards": [
                {"icon": "🔒", "title": "1. Zero-Trust Sanitizer", "text": "Replaces PII with synthetic tokens ([USER_1], [INTERNAL_IP_1]) in local RAM before any network transmission."},
                {"icon": "🤖", "title": "2. 3-Tier System-Level MoE", "text": "Triages 90% routine alerts offline on local workstation GPUs for $0 cost, cascading to cloud models only when needed."},
                {"icon": "🔒", "title": "3. AST Sandbox Guard", "text": "Parses Python AST syntax trees (ast.parse) to block dangerous shell calls (os.system) before execution."},
                {"icon": "📄", "title": "4. Courtroom PDF Briefs", "text": "Generates 1-page executive courtroom-ready incident briefs in < 30 seconds with dual-view evidence records."}
            ],
            "notes": "SENTINEL solves every drawback: A Zero-Trust Sanitizer isolates PII; a 3-Tier AI Router cuts software costs by 85%; an AST Sandbox blocks code injection; and an executive PDF report is generated in under 30 seconds."
        },
        # SLIDE 5: FULL END-TO-END SYSTEM ARCHITECTURE & FLOWCHART MAPPING
        {
            "takeaway": "ARCHITECTURE FLOW: Telemetry ──► Level 1 Sanitizer ──► Level 2 MITRE ──► Level 3 MoE Router ──► Level 5 AST ──► PDF",
            "slide_title": "🏗️ Full System Architecture & Flowchart Mapping",
            "subtitle": "Step-by-Step Data Flow from Ingestion to Courtroom Report",
            "cards": [
                {"icon": "1️⃣", "title": "Ingestion & Sanitization", "text": "• Raw SIEM Log ──► Level 1 Data Sanitizer\n• Scrubs PII to [USER_1] in local RAM.\n• Prompt Injection Firewall neutralizes attacks."},
                {"icon": "2️⃣", "title": "MITRE & 3-Tier AI Routing", "text": "• Level 2 MITRE Mapper tags T1110.\n• Level 3 Router: Tier 1 Local GPU ($0) ──► Tier 2 Groq/Gemini ──► Tier 3 550B."},
                {"icon": "3️⃣", "title": "Graph & AST Sandbox", "text": "• Level 4 Attack Graph Builder maps edges.\n• Level 5 AST Sandbox inspects syntax trees & blocks os.system()."},
                {"icon": "4️⃣", "title": "RAG, HITL & PDF Brief", "text": "• Level 6 ChromaDB RAG stores vectors.\n• Level 7/8 Dual-View HITL Officer Gate.\n• Level 9/10 Courtroom PDF generated."}
            ],
            "notes": "Judges, here is our full end-to-end architecture flowchart: Raw telemetry passes through Level 1 PII Sanitization, Level 2 MITRE Mapping, Level 3 3-Tier AI Routing, Level 5 AST Security Inspection, and Level 10 Courtroom PDF Generation."
        },
        # SLIDE 6: ARCHITECTURAL PROVENANCE & OPEN-SOURCE ADAPTATION
        {
            "takeaway": "PROVENANCE: SENTINEL adapts industry open-source standards, fixing critical security, privacy, and cost flaws.",
            "slide_title": "🏛️ Architectural Provenance & Open-Source Adaptations",
            "subtitle": "How SENTINEL Improves Upon Established Open-Source Projects",
            "cards": [
                {"icon": "🔒", "title": "Sanitizer (from Presidio)", "text": "<b>Presidio Flaw</b>: Destroys PII or sends raw text.\n<b>SENTINEL Fix</b>: Reversible token mapping ([USER_1]) stored in local RAM."},
                {"icon": "🤖", "title": "AI Router (from RouteLLM)", "text": "<b>RouteLLM Flaw</b>: Leaks prompt metadata to cloud.\n<b>SENTINEL Fix</b>: Sanitizes prompts first; processes 90% offline on local GPU ($0)."},
                {"icon": "🔒", "title": "AST Guard (from PyInquirer)", "text": "<b>PyInquirer Flaw</b>: Unsafe eval() runtime wrappers.\n<b>SENTINEL Fix</b>: AST syntax tree parsing (ast.parse) blocking os.system()."},
                {"icon": "🧠", "title": "RAG Memory (from ChromaDB)", "text": "<b>LangChain Flaw</b>: Stores raw PII in vectors.\n<b>SENTINEL Fix</b>: Embeds only PII-scrubbed threat vectors in ChromaDB RAG."}
            ],
            "notes": "We adapted established open-source projects—Microsoft Presidio for sanitization, RouteLLM for routing, and ChromaDB for vector memory—while fixing their security and privacy flaws."
        },
        # SLIDE 7: ZERO-TRUST SANITIZER & DUAL-VIEW INTERFACE
        {
            "takeaway": "DUAL-VIEW ISOLATION: Cloud AI only sees synthetic tokens; authorized officers unmask identities locally.",
            "slide_title": "🔒 Zero-Trust Sanitizer & Dual-View Interface",
            "subtitle": "Complete PII Cloud Isolation + Authorized Police Re-Identification",
            "code_box": "❌ RAW LOG (Local Workstation Only):\n\"Failed SSH login for officer.sharma@keralapolice.gov.in from 192.168.1.45 on port 22.\"\n\n✅ [Cloud / AI View] (Sent to Groq / Gemini / OpenRouter):\n\"Failed SSH login for [USER_1] from [INTERNAL_IP_1] on port 22.\"\n\n🔑 [Officer Re-Identified View] (Authorized Police Officer Only):\n\"Failed SSH login for officer.sharma@keralapolice.gov.in from 192.168.1.45 on port 22.\"",
            "notes": "Cloud AI engines only ever see PII-free synthetic tokens like [USER_1] from [INTERNAL_IP_1]. The unmasking key lives strictly inside local RAM, accessible only by authorized officers with role tokens."
        },
        # SLIDE 8: 3-TIER SYSTEM-LEVEL MoE AI ROUTING ENGINE
        {
            "takeaway": "AI ROUTING: 90% routine alerts triaged offline on GPU ($0 cost); complex zero-days cascade to 550B cloud models.",
            "slide_title": "🤖 3-Tier System-Level MoE AI Routing Engine",
            "subtitle": "85%+ Software Cost Optimization + Automatic Cascade Failover",
            "cards": [
                {"icon": "🖥️", "title": "Tier 1: Local GPU Ollama", "text": "deepseek-r1:8b / llama3.2:1b\n100% Offline GPU AI execution ($0.00 cost, 90% routine triage)."},
                {"icon": "⚡", "title": "Tier 2: Groq & Gemini Flash", "text": "DeepSeek 70B @ 300 t/s & Gemini Flash 2M Context Window\nUltra-fast reasoning & massive log file ingestion."},
                {"icon": "🌌", "title": "Tier 3: OpenRouter Free 550B", "text": "nvidia/nemotron-3-ultra-550b-a55b:free\n550 Billion Parameter intelligence for zero-day threat analysis."}
            ],
            "notes": "Our 3-Tier Router processes 90% of routine alerts locally on GPU for $0 cost. For zero-day threats, SENTINEL cascades to Groq 70B, Gemini 2M Context, or OpenRouter 550B models."
        },
        # SLIDE 9: AST CODE SANDBOX GUARD & ACTIVE DEFENSE
        {
            "takeaway": "CODE SAFETY: AST syntax tree inspection guarantees zero execution of malicious shell commands.",
            "slide_title": "🔒 AST Safe AI Code Execution Sandbox Guard",
            "subtitle": "Syntax Tree Inspection Blocking Command Injection",
            "code_box": "AI Code Input ──► ast.parse() ──► ASTSecurityVisitor Inspection\n\n✅ SAFE CODE: base64.b64decode(\"aGVsbG8=\") ──► EXECUTED IN RESTRICTED NAMESPACE\n❌ MALICIOUS: os.system(\"rm -rf /\")         ──► BLOCKED INSTANTLY (AST Security Violation)",
            "notes": "When AI generates Python scripts to de-obfuscate malware payloads, SENTINEL inspects the Python AST syntax tree first. If dangerous calls like os.system() are detected, SENTINEL blocks them instantly."
        },
        # SLIDE 10: COURTROOM PDF BRIEFS & IMMUTABLE AUDIT TRAIL
        {
            "takeaway": "EVIDENCE INTEGRITY: Append-only audit logging and courtroom-ready 1-page PDF briefs generated in < 30s.",
            "slide_title": "📜 Courtroom PDF Briefs & Immutable Audit Trail",
            "subtitle": "Courtroom-Ready Reports Generated in < 30 Seconds",
            "cards": [
                {"icon": "📄", "title": "Courtroom PDF Briefs", "text": "Generates 1-page executive PDF incident briefs using ReportLab for law enforcement and judicial review."},
                {"icon": "📜", "title": "Immutable Audit Trail", "text": "Writes append-only JSON logs (data/audit/sentinel_audit_trail.jsonl) with zero identity_map exposure."}
            ],
            "notes": "SENTINEL logs every triage event to an append-only audit trail and generates a 1-page courtroom-ready executive PDF report in under 30 seconds."
        },
        # SLIDE 11: ENTERPRISE IMPACT & VERIFICATION VICTORY
        {
            "takeaway": "VERIFICATION: SENTINEL is 100% verified operational across all 10 architectural levels.",
            "slide_title": "🏆 Enterprise Impact & System Verification",
            "subtitle": "All 10 Architectural Levels Verified Operational",
            "cards": [
                {"icon": "🟢", "title": "Zero-Trust PII Isolation", "text": "Competitors leak police PII; SENTINEL is 100% privacy-compliant."},
                {"icon": "🟢", "title": "85%+ Cost Optimization", "text": "Competitors burn cloud API fees; SENTINEL runs local GPU AI ($0)."},
                {"icon": "🟢", "title": "AST Code Safety", "text": "Competitors risk command injection; SENTINEL enforces syntax safety."},
                {"icon": "🟢", "title": "10/10 Passed", "text": "All 10 architectural levels verified operational and live on GitHub master branch."}
            ],
            "notes": "To conclude, judges: SENTINEL delivers Zero-Trust Privacy, 3-Tier AI Cost Optimization, AST Code Security, and Courtroom PDF Briefs. All 10 architectural levels are verified live. Thank you!"
        }
    ]

    for data in slides_content:
        slide = prs.slides.add_slide(blank_layout)

        # Background Fill
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_NAVY

        # Slide Header Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.4))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True

        p_t = tf_t.paragraphs[0]
        p_t.text = data["slide_title"]
        p_t.font.name = "Helvetica"
        p_t.font.bold = True
        p_t.font.size = Pt(24)
        p_t.font.color.rgb = TEXT_WHITE

        p_sub = tf_t.add_paragraph()
        p_sub.text = data.get("subtitle", "")
        p_sub.font.name = "Helvetica"
        p_sub.font.bold = True
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = ACCENT_CYAN
        p_sub.space_before = Pt(3)

        if "takeaway" in data:
            p_take = tf_t.add_paragraph()
            p_take.text = f"📌 {data['takeaway']}"
            p_take.font.name = "Helvetica"
            p_take.font.bold = True
            p_take.font.size = Pt(11)
            p_take.font.color.rgb = ACCENT_YELLOW
            p_take.space_before = Pt(4)

        # Render Visual Cards or Code Box
        if "code_box" in data:
            code_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
            tf_c = code_box.text_frame
            tf_c.word_wrap = True
            p_code = tf_c.paragraphs[0]
            p_code.text = data["code_box"]
            p_code.font.name = "Courier New"
            p_code.font.size = Pt(14)
            p_code.font.color.rgb = ACCENT_BLUE
        else:
            cards = data.get("cards", [])
            num_cards = len(cards)
            card_width = Inches(2.7) if num_cards == 4 else (Inches(3.6) if num_cards == 3 else Inches(5.5))
            gap = Inches(0.3)
            start_left = Inches(0.8)

            for i, c_data in enumerate(cards):
                left = start_left + i * (card_width + gap) if num_cards <= 4 else start_left + (i % 2) * (Inches(5.6) + gap)
                top = Inches(2.3) if i < 4 else Inches(4.6)

                card_box = slide.shapes.add_textbox(left, top, card_width, Inches(2.1))
                tf_card = card_box.text_frame
                tf_card.word_wrap = True

                p_c_title = tf_card.paragraphs[0]
                p_c_title.text = f"{c_data['icon']} {c_data['title']}"
                p_c_title.font.name = "Helvetica"
                p_c_title.font.bold = True
                p_c_title.font.size = Pt(14)
                p_c_title.font.color.rgb = ACCENT_YELLOW
                p_c_title.space_after = Pt(4)

                p_c_text = tf_card.add_paragraph()
                p_c_text.text = c_data["text"]
                p_c_text.font.name = "Helvetica"
                p_c_text.font.size = Pt(10.5)
                p_c_text.font.color.rgb = TEXT_MUTED

        # Speaker Notes Section
        notes_slide = slide.notes_slide
        tf_notes = notes_slide.notes_text_frame
        tf_notes.text = data.get("notes", "")

    prs.save(output_path)
    print(f"🎉 Master Structured PowerPoint Deck created successfully at: {output_path}")
    return output_path

if __name__ == "__main__":
    create_master_structured_pptx()
