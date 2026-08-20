import os
import sys
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

# --- COLOR PALETTE (Professional White Background Theme) ---
BG_COLOR = RGBColor(255, 255, 255)         # Pure White
NAVY_HEADER = RGBColor(10, 25, 47)         # Deep Navy (#0A192F)
TEXT_DARK = RGBColor(30, 41, 59)          # Slate Dark (#1E293B)
TEXT_MUTED = RGBColor(71, 85, 105)        # Muted Gray (#475569)
BLUE_ACCENT = RGBColor(37, 99, 235)       # Royal Blue (#2563EB)
CARD_BG = RGBColor(248, 250, 252)         # Off-White Card (#F8FAFC)
CARD_BORDER = RGBColor(226, 232, 240)     # Light Gray Border (#E2E8F0)
WHITE = RGBColor(255, 255, 255)

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, title_text, category_text="SHIELD AI | TEAM LEAD PRESENTATION"):
    # Category / Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_tag = tag_box.text_frame
    tf_tag.word_wrap = True
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = category_text.upper()
    p_tag.font.name = "Calibri"
    p_tag.font.size = Pt(10)
    p_tag.font.bold = True
    p_tag.font.color.rgb = BLUE_ACCENT

    # Main Title Header
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Calibri"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = NAVY_HEADER

    # Header Divider Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE_ACCENT
    line.line.color.rgb = BLUE_ACCENT

def add_footer(slide, current_slide, total_slides=10):
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.3))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Sri Sai Ram Engineering College | Department of CSE | SHIELD AI Presentation — Slide {current_slide} of {total_slides}"
    p.font.name = "Calibri"
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED

def add_card(slide, left, top, width, height, title, content_bullets, title_color=NAVY_HEADER):
    # Base rounded card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(1)

    # Content Frame
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), height - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True

    p_header = tf.paragraphs[0]
    p_header.text = title
    p_header.font.name = "Calibri"
    p_header.font.size = Pt(14)
    p_header.font.bold = True
    p_header.font.color.rgb = title_color
    p_header.space_after = Pt(8)

    for bullet in content_bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.name = "Calibri"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    apply_background(s1)

    # Title Card Accent Box
    accent_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(0.15), Inches(4.8))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = BLUE_ACCENT
    accent_bar.line.color.rgb = BLUE_ACCENT

    # Main Title
    tb1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(11.0), Inches(2.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "SHIELD AI"
    p1.font.name = "Calibri"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = NAVY_HEADER

    p2 = tf1.add_paragraph()
    p2.text = "Autonomous Cyber Defence & Security Intelligence Platform"
    p2.font.name = "Calibri"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = BLUE_ACCENT
    p2.space_before = Pt(6)

    p3 = tf1.add_paragraph()
    p3.text = "Privacy-Preserving SOC Alert Triage, Multi-Tier AI Routing & Automated Threat Investigation Engine"
    p3.font.name = "Calibri"
    p3.font.size = Pt(13)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(8)

    # Team Members Metadata Box
    card_meta = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(3.8), Inches(11.3), Inches(2.3))
    card_meta.fill.solid()
    card_meta.fill.fore_color.rgb = CARD_BG
    card_meta.line.color.rgb = CARD_BORDER

    tb_meta = s1.shapes.add_textbox(Inches(1.4), Inches(3.9), Inches(10.9), Inches(2.1))
    tf_m = tb_meta.text_frame
    tf_m.word_wrap = True

    p_m0 = tf_m.paragraphs[0]
    p_m0.text = "PROJECT LEADERSHIP & TEAM COMPOSITION"
    p_m0.font.name = "Calibri"
    p_m0.font.size = Pt(13)
    p_m0.font.bold = True
    p_m0.font.color.rgb = NAVY_HEADER
    p_m0.space_after = Pt(6)

    members_text = [
        "• Team Lead (TL): Gokula Kannan M (SEC25CS196 | CSE-A) — Systems Architecture & Venture Strategy",
        "• Team Member 1 (M1): Lakshan M (SEC25CS036 | CSE-A) — Full-Stack Dashboard Architect & HITL UI",
        "• Team Member 2 (M2): Sivabalan T (SEC25CS101 | CSE-C) — AI Engine Architecture & Zero-Trust Sanitizer",
        "• Faculty Supervisor: Dr. A. SHEELA (Associate Professor, Dept. of Computer Science & Engineering)",
        "• Institution: Sri Sai Ram Engineering College, Chennai | Academic Year: 2026–2027"
    ]
    for m in members_text:
        p = tf_m.add_paragraph()
        p.text = m
        p.font.name = "Calibri"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(3)

    add_footer(s1, 1)

    # ==========================================
    # SLIDE 2: EXECUTIVE SUMMARY & VENTURE OVERVIEW
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    apply_background(s2)
    add_header(s2, "Executive Summary & Team Lead Perspective")

    add_card(s2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1), 
             "The Problem We Solved", [
                 "• Severe Alert Fatigue: Enterprise SOCs ingest over 5,000 logs daily, with manual investigation taking 30–45 minutes per alert.",
                 "• Unreviewed Threat Exposure: Over 70% of security alerts go unexamined due to human analyst capacity limits.",
                 "• Cloud Telemetry Privacy Leakage: Transmitting raw logs to public cloud AI violates DPDP Act 2023 and risks exposing internal IP maps, emails, and credentials.",
                 "• Air-Gapped Operational Barrier: Military labs and police cyber cells operate without internet, making cloud AI non-viable."
             ])

    add_card(s2, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.1), 
             "Our Architectural Breakthrough", [
                 "• Zero-Trust Local Data Sanitizer: Redacts 100% of PII, IPs, MACs, and tokens into synthetic handles ([USER_1], [INTERNAL_IP_1]) in local RAM.",
                 "• Prompt Injection Firewall: Neutralizes adversarial prompt overrides in raw logs before AI processing.",
                 "• 3-Tier Hybrid AI Model Router: Resolves ~90% of routine alerts 100% offline on Tier-1 Ollama (deepseek-r1:8b) at $0 compute cost.",
                 "• Human-in-the-Loop Active Defense: Enforces strict analyst click-approval before executing containment actions."
             ], title_color=BLUE_ACCENT)

    add_footer(s2, 2)

    # ==========================================
    # SLIDE 3: PROBLEM CONTEXT & SOC BOTTLENECKS
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    apply_background(s3)
    add_header(s3, "Problem Context & The SOC Bottleneck")

    add_card(s3, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.1),
             "1. Triage Bottleneck", [
                 "• 5,000+ Daily SIEM Alerts",
                 "• 30–45 Minutes per Manual Triage",
                 "• 70%+ Alerts Left Unreviewed",
                 "• High Analyst Burnout & Turnover",
                 "• Extended Attacker Dwell Time"
             ])

    add_card(s3, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.1),
             "2. Data Privacy & Compliance", [
                 "• Cloud Egress Exposure Risks",
                 "• Internal IP Topographies Leakage",
                 "• Staff Credentials & PII Compromise",
                 "• DPDP Act 2023 Non-Compliance",
                 "• Forensic Chain-of-Custody Breach"
             ])

    add_card(s3, Inches(8.8), Inches(1.6), Inches(3.733), Inches(5.1),
             "3. Air-Gapped Barrier", [
                 "• Zero Internet Egress Mandate",
                 "• High Commercial Cloud API Costs",
                 "• Unreliable External Model APIs",
                 "• Lack of Edge Hardware Support",
                 "• Complex Deployment Pipelines"
             ])

    add_footer(s3, 3)

    # ==========================================
    # SLIDE 4: PROPOSED SYSTEM & CORE VALUE PROPOSITION
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    apply_background(s4)
    add_header(s4, "Proposed System & Core Value Proposition")

    add_card(s4, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1),
             "Key Technical Innovations", [
                 "1. Async FastAPI Webhook Receiver: Non-blocking ingestion bridge receiving Wazuh SIEM JSON streams.",
                 "2. In-RAM Regex Data Sanitizer: Redacts emails, IPv4/v6, MACs, and API tokens while storing lookup maps in RAM.",
                 "3. Prompt-Injection Neutralizer: Replaces malicious prompt overrides with [NEUTRALIZED_PROMPT_INJECTION] tokens.",
                 "4. AST Python Code Execution Sandbox: Safely parses script syntax trees using ast.parse() to block dangerous calls."
             ])

    add_card(s4, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.1),
             "Measurable System Value", [
                 "• Machine-Speed Triage: Resolves alert investigation and builds reports in <30 seconds per incident.",
                 "• 100% Data Sovereignty: Sensitive identifiers never cross external network boundaries.",
                 "• $0 Operational Base Cost: Resolves ~90% of routine alerts locally on GPU hardware via Ollama.",
                 "• Courtroom-Ready Evidence: Generates verifiable PDF incident briefs backed by JSONL audit logs."
             ], title_color=BLUE_ACCENT)

    add_footer(s4, 4)

    # ==========================================
    # SLIDE 5: ARCHITECTURAL FRAMEWORK & 3-TIER ROUTER
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    apply_background(s5)
    add_header(s5, "3-Tier System Architecture & Model Router")

    add_card(s5, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.1),
             "Tier 1: Local Offline GPU", [
                 "• Runtime: Ollama Local Inference",
                 "• Model: deepseek-r1:8b / llama3.2:1b",
                 "• Target: ~90% Routine Alerts",
                 "• Cost: $0 Marginal Software Cost",
                 "• Network: 100% Air-Gapped Offline"
             ])

    add_card(s5, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.1),
             "Tier 2: Fast Cloud Reasoning", [
                 "• Endpoint: Groq Cloud API",
                 "• Model: deepseek-r1-distill-llama-70b",
                 "• Trigger: Complex Multi-Stage Threats",
                 "• Payload: Anonymized Tokens Only",
                 "• Speed: 300+ tokens/second"
             ])

    add_card(s5, Inches(8.8), Inches(1.6), Inches(3.733), Inches(5.1),
             "Tier 3: Enterprise Multi-Modal", [
                 "• Endpoint: Gemini / OpenRouter / OpenAI",
                 "• Model: Gemini 2.0 Flash / GPT-4o",
                 "• Trigger: Multi-Gigabyte Disk Dumps",
                 "• Context: Up to 2M Token Context",
                 "• Control: Policy-Governed Escalation"
             ])

    add_footer(s5, 5)

    # ==========================================
    # SLIDE 6: INTELLIGENCE & THREAT ANALYSIS ENGINES
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    apply_background(s6)
    add_header(s6, "Intelligence, Vector Memory & Attack Graph Correlation")

    add_card(s6, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1),
             "Threat Correlation & RAG Memory", [
                 "• MITRE ATT&CK Mapping: Rule-assisted classifier mapping log features to tactics (Credential Access) & techniques (T1110 Brute Force).",
                 "• Vector Threat Memory: Embedded ChromaDB database storing dense embeddings of past security incidents.",
                 "• Cosine Similarity Retrieval: Queries historical incident resolutions to provide contextual precedent for new alerts.",
                 "• Automated Pattern Matching: Reduces investigation overhead for recurring attack vectors."
             ])

    add_card(s6, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.1),
             "Attack Graph Synthesis (DAG)", [
                 "• Temporal Entity Correlation: Clusters events occurring across sliding time windows and shared entity tokens.",
                 "• Intrusion Lifecycle Mapping: Reconstructs sequential attack steps from Initial Access (T1078) to Exfiltration (T1041).",
                 "• Directed Acyclic Graphs (DAG): Generates visual attack graph nodes and edges for analyst review.",
                 "• Root-Cause Discovery: Exposes hidden low-and-slow intruder lateral movement."
             ], title_color=BLUE_ACCENT)

    add_footer(s6, 6)

    # ==========================================
    # SLIDE 7: SAFETY, GOVERNANCE & RESPONSE MECHANISMS
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    apply_background(s7)
    add_header(s7, "Safety, Governance & Active Defense Mechanisms")

    add_card(s7, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.1),
             "AST Code Sandbox", [
                 "• Python ast.parse() Syntax Tree",
                 "• Blocks Dangerous Primitives",
                 "• Filters os, sys, subprocess, eval",
                 "• Safe Payload De-obfuscation",
                 "• Prevents Arbitrary Code Execution"
             ])

    add_card(s7, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.1),
             "HITL Authorization Gate", [
                 "• Human-in-the-Loop Analyst Gate",
                 "• Role-Based Access Control (RBAC)",
                 "• Mandatory Click-Approval Modal",
                 "• Prevents AI Action Hallucination",
                 "• Officer Authorization Check"
             ])

    add_card(s7, Inches(8.8), Inches(1.6), Inches(3.733), Inches(5.1),
             "Mock Active Response", [
                 "• Simulated Firewall IP Blocking",
                 "• Process Termination Controller",
                 "• Host Isolation Simulation",
                 "• Non-Disruptive Testing Mode",
                 "• Production Allowlist Safety"
             ])

    add_footer(s7, 7)

    # ==========================================
    # SLIDE 8: TEAM COMPOSITION & ROLE MATRIX
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    apply_background(s8)
    add_header(s8, "Team Composition & Division of Responsibilities")

    # Team Members Table
    table_shape = s8.shapes.add_table(5, 4, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.1))
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(3.5)
    table.columns[3].width = Inches(4.033)

    headers = ["Team Member", "Role & Section", "Core Technical Domain", "Project Deliverables & Contributions"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_HEADER
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Calibri"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    rows_data = [
        ("Gokula Kannan M", "Team Lead (TL)\nSEC25CS196 | CSE-A", "Systems Architecture, Asynchronous Backends & Venture Strategy", "Overall system architecture, 3-Tier AI Router design, sprint management, leading national pitch at Hac'KP 2026 @ Zoho."),
        ("Lakshan M", "Team Member 1 (M1)\nSEC25CS036 | CSE-A", "Full-Stack Web Architecture, React.js UI & WebSockets", "Real-time incident dashboard engineering, WebSockets live feed, and Human-in-the-Loop (HITL) analyst approval interface."),
        ("Sivabalan T", "Team Member 2 (M2)\nSEC25CS101 | CSE-C", "AI Engine Architecture, Zero-Trust Privacy & Quantization", "Development of Zero-Trust Data Sanitizer (src/sanitizer.py), 3-Tier AI Router (src/router.py), and local GGUF quantization."),
        ("Dr. A. SHEELA", "Faculty Supervisor\nAssoc. Prof., CSE", "Academic Supervision, Research Guidance & Governance", "Project oversight, academic literature review guidance, and regulatory compliance alignment for Sairam Institutions.")
    ]

    for r_idx, row in enumerate(rows_data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Calibri"
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_DARK

    add_footer(s8, 8)

    # ==========================================
    # SLIDE 9: IEEE LITERATURE MATRIX & SDGS
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    apply_background(s9)
    add_header(s9, "IEEE Literature Survey & UN SDG Alignment")

    add_card(s9, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1),
             "IEEE Academic Literature Matrix (6 Papers)", [
                 "• Vasilev et al. (IEEE 2026): Proves 8B local LLMs reduce false positives by 85% — validates Tier-1 Ollama.",
                 "• Zhang & Liu (IEEE TIFS 2025): Proves tokenized logs retain 100% semantic utility — validates DataSanitizer.",
                 "• Nguyen & Pham (IEEE Access 2025): Proves 91.4% precision in rule mapping — validates mitre_mapper.py.",
                 "• Kumar et al. (IEEE EMBC 2025): Proves model cascading cuts costs by 78% — validates 3-Tier AI Router.",
                 "• Bansal et al. (IEEE SPW 2024): Proves structured briefs reduce MTTT by 64% — validates PDF exporter."
             ])

    add_card(s9, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.1),
             "UN Sustainable Development Goals (SDGs)", [
                 "• SDG 9 (Industry, Innovation & Infrastructure): Builds resilient public cyber defense infrastructure with zero cloud dependency.",
                 "• SDG 16 (Peace, Justice & Strong Institutions): Enforces data privacy, evidence chain-of-custody, and court-admissible audit logs.",
                 "• SDG 8 (Decent Work & Economic Growth): Automates routine Tier-1 triage to eliminate human analyst burnout.",
                 "• SDG 12 (Responsible Consumption & Production): Reduces cloud energy footprints via local edge GPU inference."
             ], title_color=BLUE_ACCENT)

    add_footer(s9, 9)

    # ==========================================
    # SLIDE 10: FUTURE ROADMAP & CONCLUSION
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    apply_background(s10)
    add_header(s10, "Future Commercialization Roadmap & Conclusion")

    add_card(s10, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1),
             "Future Development & Hardening Roadmap", [
                 "1. Labeled Benchmark Dataset: Benchmarking against CIC-IDS-2017 & production Wazuh SIEM telemetry.",
                 "2. Local NER Sanitizer: Integrating lightweight local ONNX NER models (spaCy) for non-regex PII entities.",
                 "3. Cryptographic Audit Log: Upgrading JSONL audit logs with HMAC hash chains & hardware TPM signing.",
                 "4. Enterprise Containerization: Packaging the full platform using Docker, Kubernetes, and CI/CD pipelines."
             ])

    add_card(s10, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.1),
             "Conclusion & Team Lead Vision", [
                 "• SHIELD AI proves that enterprise-grade SOC triage can be automated at machine speed (<30s MTTR) without compromising data privacy or incurring cloud software costs.",
                 "• By unifying Zero-Trust PII Tokenization, 3-Tier AI Routing, Vector Threat Memory, and Human-in-the-Loop Governance, SHIELD AI delivers a legally robust defense platform.",
                 "• On behalf of Team Lead Gokula Kannan M & Team SENTINEL, we thank Sairam Institutions for their continuous guidance!"
             ], title_color=BLUE_ACCENT)

    add_footer(s10, 10)

    # Save outputs
    output_path = r"c:\Users\siva2\Projects\SENTINEL\docs\SHIELD_AI_Master_Presentation.pptx"
    desktop_path = r"C:\Users\siva2\OneDrive\Desktop\SHIELD_AI_Master_Presentation.pptx"

    prs.save(output_path)
    shutil.copy2(output_path, desktop_path)
    print(f"🎉 MASTER POWERPOINT PRESENTATION CREATED: {desktop_path}")
    return desktop_path

if __name__ == "__main__":
    create_presentation()
